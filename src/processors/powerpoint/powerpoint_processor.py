"""
PowerPoint Processor - UPDATED: Now works with AccessibilityOrderExtractorV2
Updated to handle tuple return format (shape, role) from new extractor
"""

from pptx import Presentation
import os
from datetime import datetime
from markitdown import MarkItDown

# Import your new extractor instead of the old one
from .accessibility_extractor_v2 import AccessibilityOrderExtractorV2
from .content_extractor import ContentExtractor
from .diagram_analyzer import DiagramAnalyzer
from .text_processor import TextProcessor
from .markdown_converter import MarkdownConverter
from .metadata_extractor import MetadataExtractor


class PowerPointProcessor:
    """
    Main PowerPoint processor implementing dual-strategy processing architecture
    with proper semantic role information flow and group handling.
    UPDATED: Now works with AccessibilityOrderExtractorV2
    """

    def __init__(self, use_accessibility_order=True):
        """Initialize the PowerPoint processor with all component dependencies."""
        self.use_accessibility_order = use_accessibility_order

        # Initialize specialized components for XML-based processing
        # UPDATED: Use new extractor class
        self.accessibility_extractor = AccessibilityOrderExtractorV2(use_accessibility_order)
        self.content_extractor = ContentExtractor()
        self.diagram_analyzer = DiagramAnalyzer()
        self.text_processor = TextProcessor()
        self.markdown_converter = MarkdownConverter()
        self.metadata_extractor = MetadataExtractor()

        # Initialize MarkItDown for fallback processing
        self.markitdown = MarkItDown()

        # Supported file format configuration
        self.supported_formats = ['.pptx', '.ppt']

    def _has_xml_access(self, file_path):
        """Check if XML-based processing is possible for the given file."""
        try:
            prs = Presentation(file_path)
            if len(prs.slides) > 0:
                # UPDATED: Simple check since new extractor handles this internally
                return True
            return False
        except Exception:
            return False

    def extract_slide_data(self, slide, slide_number):
        """
        Extract content from individual slide using coordinated component pipeline.
        UPDATED: Handle tuple return format from AccessibilityOrderExtractorV2
        """
        print(f"\n=== PROCESSING SLIDE {slide_number} ===")

        # UPDATED: Get shapes with roles from new extractor
        shapes_with_roles = self.accessibility_extractor.get_slide_reading_order(slide, slide_number)

        print(f"DEBUG: New extractor returned {len(shapes_with_roles)} shape-role pairs")

        slide_data = {
            "slide_number": slide_number,
            "content_blocks": [],
            "extraction_method": "semantic_accessibility_order_v2"  # Updated method name
        }

        # UPDATED: The new extractor doesn't expand groups at slide level
        # Groups are handled during recursive expansion within the extractor
        groups_were_expanded = True  # Always true for new extractor

        print(f"DEBUG: Groups were expanded: {groups_were_expanded}")

        # Show what shapes we're processing
        print(f"DEBUG: Shapes to process:")
        for i, (shape, role) in enumerate(shapes_with_roles[:10]):  # Show first 10
            shape_type = str(shape.shape_type).split('.')[-1] if hasattr(shape.shape_type, '__str__') else 'unknown'
            try:
                text_preview = ""
                if hasattr(shape, 'text') and shape.text:
                    text_preview = shape.text.strip()[:30] + "..."
                elif hasattr(shape, 'text_frame') and shape.text_frame:
                    text_preview = shape.text_frame.text.strip()[:30] + "..."
                else:
                    text_preview = "No text"
                print(f"  {i + 1}. {shape_type} ({role}): {text_preview}")
            except:
                print(f"  {i + 1}. {shape_type} ({role}): Error getting text")

        # UPDATED: Extract content from each shape using pre-determined semantic role
        processed_count = 0
        for i, (shape, semantic_role) in enumerate(shapes_with_roles):
            print(f"\nDEBUG: Processing shape {i + 1}/{len(shapes_with_roles)} with role '{semantic_role}'")

            # UPDATED: Pass the semantic role directly to content extractor
            block = self.content_extractor.extract_shape_content(
                shape,
                self.text_processor,
                self.accessibility_extractor,
                groups_already_expanded=groups_were_expanded,
                semantic_role=semantic_role  # NEW: Pass pre-determined role
            )
            if block:
                slide_data["content_blocks"].append(block)
                processed_count += 1
                print(f"DEBUG: Added content block (type: {block.get('type', 'unknown')}, role: {semantic_role})")
            else:
                print(f"DEBUG: Shape produced no content block")

        print(
            f"DEBUG: Slide {slide_number} final result: {processed_count} content blocks from {len(shapes_with_roles)} shapes")
        return slide_data

    def debug_accessibility_order(self, file_path, slide_number=1):
        """Debug method for analyzing reading order extraction and semantic roles."""
        try:
            if not self._has_xml_access(file_path):
                print(f"❌ XML not available for {file_path}")
                print("Would use MarkItDown fallback in production")
                return

            prs = Presentation(file_path)
            if slide_number > len(prs.slides):
                print(f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides.")
                return

            slide = prs.slides[slide_number - 1]
            print(f"🎯 XML available - debugging sophisticated processing with semantic roles...")

            print(f"\n=== DEBUGGING SLIDE {slide_number} READING ORDER WITH SEMANTIC ROLES ===")
            print(f"Total shapes on slide: {len(slide.shapes)}")

            # UPDATED: Test new extractor directly
            shapes_with_roles = self.accessibility_extractor.get_slide_reading_order(slide, slide_number)
            print(f"✅ New extractor returned: {len(shapes_with_roles)} shape-role pairs")

            # Show shape information with semantic roles
            print("\n🎯 SHAPE ORDER WITH SEMANTIC ROLES:")
            for i, (shape, semantic_role) in enumerate(shapes_with_roles):
                shape_type = str(shape.shape_type).split('.')[-1]

                text_preview = ""
                try:
                    if hasattr(shape, 'text') and shape.text:
                        text_preview = shape.text.strip()[:40] + "..."
                    elif hasattr(shape, 'text_frame') and shape.text_frame:
                        text_preview = shape.text_frame.text.strip()[:40] + "..."
                except:
                    text_preview = "No text"

                print(f"  {i + 1}. [{shape_type}] SEMANTIC_ROLE: {semantic_role} | {text_preview}")

            # Test content extraction
            print(f"\n🎯 TESTING CONTENT EXTRACTION:")
            slide_data = self.extract_slide_data(slide, slide_number)
            print(f"✅ Content blocks extracted: {len(slide_data['content_blocks'])}")

            for i, block in enumerate(slide_data['content_blocks']):
                block_type = block.get('type', 'unknown')
                semantic_role = block.get('semantic_role', 'unknown')

                if block_type == 'text' and block.get('paragraphs'):
                    text_preview = block['paragraphs'][0].get('clean_text', '')[:40] + "..."
                elif block_type == 'group':
                    child_count = len(block.get('extracted_blocks', []))
                    text_preview = f"Group with {child_count} children"
                else:
                    text_preview = f"{block_type} content"

                print(f"  {i + 1}. [{block_type}] SEMANTIC: {semantic_role} | {text_preview}")

        except Exception as e:
            print(f"Debug failed: {str(e)}")
            import traceback
            traceback.print_exc()

    def configure_extraction_method(self, use_accessibility_order):
        """Configure reading order extraction method for all processing."""
        self.use_accessibility_order = use_accessibility_order
        self.accessibility_extractor.accessibility_order = use_accessibility_order  # UPDATED: Different attribute name

    # ... rest of the methods remain the same ...
    def convert_pptx_to_markdown_enhanced(self, file_path, convert_slide_titles=True):
        """
        Main entry point implementing XML-first processing with MarkItDown fallback.
        NOTE: convert_slide_titles parameter kept for compatibility but XML semantic roles now control titles.
        """
        try:
            # Determine processing strategy based on XML accessibility
            if self._has_xml_access(file_path):
                return self._sophisticated_xml_processing(file_path, convert_slide_titles)
            else:
                return self._simple_markitdown_processing(file_path)

        except Exception as e:
            raise Exception(f"Error processing PowerPoint file: {str(e)}")

    def _sophisticated_xml_processing(self, file_path, convert_slide_titles):
        """
        Execute full-featured processing pipeline when XML is accessible.
        Updated to ensure semantic role information flows through properly.
        """
        print("🎯 Using sophisticated XML-based processing with semantic roles...")

        # Load presentation for full processing
        prs = Presentation(file_path)

        # Extract comprehensive metadata for document context
        pptx_metadata = self.metadata_extractor.extract_pptx_metadata(prs, file_path)

        # Process entire presentation through component pipeline
        structured_data = self.extract_presentation_data(prs)

        # Convert structured data to clean markdown with semantic role awareness
        markdown = self.markdown_converter.convert_structured_data_to_markdown(
            structured_data, convert_slide_titles=False  # XML controls titles now
        )

        # Enhance with metadata context for Claude AI processing
        markdown_with_metadata = self.metadata_extractor.add_pptx_metadata_for_claude(
            markdown, pptx_metadata
        )

        # ENHANCED: Use direct slide analysis for comprehensive diagram detection
        print("🎯 Running enhanced diagram analysis with direct slide access...")
        diagram_analysis = self.diagram_analyzer.analyze_slides_for_diagrams(
            slides=list(prs.slides),  # Pass raw slides for comprehensive analysis
            structured_data=structured_data  # Fallback if needed
        )
        if diagram_analysis:
            print("✅ Enhanced diagram analysis found potential diagrams")
            markdown_with_metadata += "\n\n" + diagram_analysis
        else:
            print("❌ Enhanced diagram analysis found no diagrams")

        return markdown_with_metadata

    def _simple_markitdown_processing(self, file_path):
        """Execute simple fallback processing using MarkItDown library."""
        print("📄 XML not available - using MarkItDown fallback...")

        try:
            result = self.markitdown.convert(file_path)

            try:
                markdown_content = result.markdown
            except AttributeError:
                try:
                    markdown_content = result.text_content
                except AttributeError:
                    raise Exception("Neither 'markdown' nor 'text_content' attribute found on result object")

            metadata_comment = f"\n<!-- Converted using MarkItDown fallback - XML not available -->\n"
            return metadata_comment + markdown_content

        except Exception as e:
            raise Exception(f"MarkItDown processing failed: {str(e)}")

    def extract_presentation_data(self, presentation):
        """
        Extract structured data from entire presentation using component coordination.
        ENHANCED: Added debugging to track slide processing.
        """
        data = {
            "total_slides": len(presentation.slides),
            "slides": []
        }

        print(f"DEBUG: Starting presentation extraction with {len(presentation.slides)} slides")

        # Process each slide individually while maintaining order
        for slide_idx, slide in enumerate(presentation.slides, 1):
            print(f"\n{'=' * 50}")
            print(f"DEBUG: EXTRACTING SLIDE {slide_idx} of {len(presentation.slides)}")
            print(f"{'=' * 50}")

            slide_data = self.extract_slide_data(slide, slide_idx)
            data["slides"].append(slide_data)

            print(f"DEBUG: Slide {slide_idx} added to presentation data")

        print(f"\nDEBUG: Presentation extraction complete - {len(data['slides'])} slides processed")
        return data

    def get_processing_summary(self, file_path):
        """Get comprehensive processing summary without performing full conversion."""
        try:
            has_xml = self._has_xml_access(file_path)

            summary = {
                "file_path": file_path,
                "has_xml_access": has_xml,
                "processing_method": "sophisticated_xml_with_semantic_roles_v2" if has_xml else "markitdown_fallback"
            }

            if has_xml:
                prs = Presentation(file_path)

                summary.update({
                    "slide_count": len(prs.slides),
                    "extraction_method": "accessibility_order_v2_with_semantic_roles",
                    "has_diagram_analysis": True,
                    "has_semantic_title_detection": True,
                    "slides_preview": []
                })

                # Generate preview for first few slides with semantic role info
                for i, slide in enumerate(prs.slides[:3], 1):
                    # UPDATED: Handle tuple return format
                    shapes_with_roles = self.accessibility_extractor.get_slide_reading_order(slide, i)

                    # Count semantic roles
                    title_count = 0
                    subtitle_count = 0
                    content_count = 0

                    for shape, semantic_role in shapes_with_roles:
                        if semantic_role == "title":
                            title_count += 1
                        elif semantic_role == "subtitle":
                            subtitle_count += 1
                        elif semantic_role == "content":
                            content_count += 1

                    slide_preview = {
                        "slide_number": i,
                        "shape_count": len(shapes_with_roles),
                        "title_shapes": title_count,
                        "subtitle_shapes": subtitle_count,
                        "content_shapes": content_count,
                        "has_text": any(hasattr(shape, 'text_frame') and shape.text_frame
                                        for shape, role in shapes_with_roles),
                        "extraction_method": "semantic_accessibility_order_v2"
                    }
                    summary["slides_preview"].append(slide_preview)
            else:
                summary.update({
                    "slide_count": "unknown",
                    "extraction_method": "markitdown_fallback",
                    "has_diagram_analysis": False,
                    "has_semantic_title_detection": False,
                    "note": "XML not available - using simple MarkItDown conversion"
                })

            return summary

        except Exception as e:
            return {"error": str(e)}


# Convenience functions for backward compatibility and simple usage
def convert_pptx_to_markdown_enhanced(file_path, convert_slide_titles=True):
    """
    Convenience function maintaining backward compatibility with simple API.
    NOTE: convert_slide_titles parameter kept for compatibility but XML semantic roles now control titles.
    """
    processor = PowerPointProcessor()
    return processor.convert_pptx_to_markdown_enhanced(file_path, convert_slide_titles)


def process_powerpoint_file(file_path, output_format="markdown", convert_slide_titles=True):
    """
    Convenience function for comprehensive file processing with multiple output options.
    NOTE: convert_slide_titles parameter kept for compatibility but XML semantic roles now control titles.
    """
    processor = PowerPointProcessor()

    if output_format == "summary":
        return processor.get_processing_summary(file_path)
    else:
        markdown_content = processor.convert_pptx_to_markdown_enhanced(file_path, convert_slide_titles)

        result = {
            "content": markdown_content,
            "format": output_format,
            "processing_method": "sophisticated_xml_with_semantic_roles_v2" if processor._has_xml_access(
                file_path) else "markitdown_fallback"
        }

        if processor._has_xml_access(file_path):
            try:
                prs = Presentation(file_path)
                result["metadata"] = processor.metadata_extractor.extract_pptx_metadata(prs, file_path)
            except Exception:
                pass

        return result