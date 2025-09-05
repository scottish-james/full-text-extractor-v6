#!/usr/bin/env python3
"""
LibreOffice Installation and Functionality Checker
Tests if LibreOffice is properly installed and can convert documents.

Usage:
    python check_libreoffice.py
    python check_libreoffice.py --test-convert sample.pptx
"""

import os
import sys
import shutil
import subprocess
import tempfile
import argparse
from pathlib import Path
import platform


class LibreOfficeChecker:
    """Check LibreOffice installation and functionality."""

    def __init__(self):
        self.libreoffice_path = None
        self.version = None
        self.system = platform.system()

    def find_libreoffice(self):
        """Find LibreOffice installation."""
        print("🔍 Searching for LibreOffice installation...")
        print(f"   System: {self.system}")

        # Platform-specific paths to check
        if self.system == "Darwin":  # macOS
            paths = [
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                "/opt/homebrew/bin/soffice",
                "/usr/local/bin/soffice",
                "/Applications/LibreOffice.app/Contents/MacOS/libreoffice"
            ]
        elif self.system == "Windows":
            paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                r"C:\Program Files\LibreOffice 7\program\soffice.exe",
                r"C:\Program Files\LibreOffice 24\program\soffice.exe"
            ]
        else:  # Linux
            paths = [
                "/usr/bin/soffice",
                "/usr/local/bin/soffice",
                "/snap/bin/libreoffice",
                "/usr/bin/libreoffice",
                "/opt/libreoffice/program/soffice",
                "/usr/lib/libreoffice/program/soffice"
            ]

        # Check each specific path
        print("\n📂 Checking common installation paths:")
        for path in paths:
            if Path(path).exists():
                print(f"   ✅ Found: {path}")
                self.libreoffice_path = path
                return path
            else:
                print(f"   ❌ Not found: {path}")

        # Check system PATH
        print("\n🔍 Checking system PATH...")
        for cmd in ["soffice", "libreoffice"]:
            path = shutil.which(cmd)
            if path:
                print(f"   ✅ Found in PATH: {path}")
                self.libreoffice_path = path
                return path

        print("   ❌ Not found in system PATH")
        return None

    def check_version(self):
        """Check LibreOffice version."""
        if not self.libreoffice_path:
            return False

        print(f"\n📋 Checking LibreOffice version...")

        try:
            result = subprocess.run(
                [self.libreoffice_path, "--version"],
                capture_output=True,
                text=True,
                timeout=10
            )

            if result.returncode == 0:
                self.version = result.stdout.strip()
                print(f"   ✅ Version: {self.version}")
                return True
            else:
                print(f"   ❌ Error getting version: {result.stderr}")
                return False

        except subprocess.TimeoutExpired:
            print("   ❌ Timeout getting version")
            return False
        except Exception as e:
            print(f"   ❌ Error: {e}")
            return False

    def test_conversion(self, test_file=None):
        """Test document conversion capability."""
        if not self.libreoffice_path:
            print("\n❌ Cannot test conversion - LibreOffice not found")
            return False

        print(f"\n🧪 Testing document conversion...")

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)

            try:
                # Create or use test file
                if test_file and Path(test_file).exists():
                    test_path = Path(test_file)
                    print(f"   Using provided file: {test_path.name}")
                else:
                    # Create a simple test document
                    test_path = temp_path / "test.html"
                    test_path.write_text("""
                    <html>
                    <body>
                        <h1>LibreOffice Test Document</h1>
                        <p>This is a test document for conversion.</p>
                        <ul>
                            <li>Item 1</li>
                            <li>Item 2</li>
                        </ul>
                    </body>
                    </html>
                    """)
                    print(f"   Created test file: {test_path.name}")

                # Test conversion to PDF
                print(f"   Converting to PDF...")

                cmd = [
                    self.libreoffice_path,
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", str(temp_path),
                    str(test_path)
                ]

                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=30
                )

                if result.returncode == 0:
                    # Check if PDF was created
                    pdf_files = list(temp_path.glob("*.pdf"))
                    if pdf_files:
                        pdf_path = pdf_files[0]
                        pdf_size = pdf_path.stat().st_size
                        print(f"   ✅ Conversion successful!")
                        print(f"   📄 Created: {pdf_path.name} ({pdf_size:,} bytes)")
                        return True
                    else:
                        print(f"   ❌ PDF not created")
                        return False
                else:
                    print(f"   ❌ Conversion failed")
                    print(f"   Error: {result.stderr}")
                    return False

            except subprocess.TimeoutExpired:
                print("   ❌ Conversion timeout")
                return False
            except Exception as e:
                print(f"   ❌ Error during conversion: {e}")
                return False

    def check_dependencies(self):
        """Check additional dependencies."""
        print(f"\n🔧 Checking additional dependencies...")

        # Check Python packages for document processing
        packages = {
            'pdf2image': 'PDF to image conversion',
            'PIL': 'Image processing (Pillow)',
            'openai': 'OpenAI Vision API'
        }

        print("\n📦 Python packages:")
        for module, description in packages.items():
            try:
                if module == 'PIL':
                    import PIL
                else:
                    __import__(module)
                print(f"   ✅ {module:<12} - {description}")
            except ImportError:
                print(f"   ❌ {module:<12} - {description}")

        # Check system tools
        print("\n🛠️ System tools:")

        # Check for Poppler (needed for pdf2image)
        if shutil.which('pdftoppm'):
            print(f"   ✅ Poppler     - PDF utilities")
        else:
            print(f"   ❌ Poppler     - PDF utilities (needed for pdf2image)")
            if self.system == "Darwin":
                print(f"      Install: brew install poppler")
            elif self.system == "Windows":
                print(f"      Download from: https://github.com/oschwartz10612/poppler-windows/releases/")
            else:
                print(f"      Install: sudo apt-get install poppler-utils")

    def print_summary(self):
        """Print installation summary and recommendations."""
        print("\n" + "=" * 50)
        print("📊 SUMMARY")
        print("=" * 50)

        if self.libreoffice_path:
            print(f"✅ LibreOffice found: {self.libreoffice_path}")
            if self.version:
                print(f"   Version: {self.version}")
        else:
            print("❌ LibreOffice NOT found")

        print("\n📚 Installation Instructions:")
        if self.system == "Darwin":  # macOS
            print("""
   macOS:
   1. Download from: https://www.libreoffice.org/download/
   2. Or install via Homebrew: brew install --cask libreoffice
   3. Make sure to allow the app in System Preferences > Security & Privacy
   """)
        elif self.system == "Windows":
            print("""
   Windows:
   1. Download from: https://www.libreoffice.org/download/
   2. Run the installer with default settings
   3. Restart your terminal/command prompt after installation
   """)
        else:  # Linux
            print("""
   Linux:
   - Ubuntu/Debian: sudo apt-get install libreoffice
   - Fedora: sudo dnf install libreoffice
   - Arch: sudo pacman -S libreoffice
   - Snap: sudo snap install libreoffice
   """)

        if not self.libreoffice_path:
            print("\n💡 After installation:")
            print("   1. Restart your terminal/IDE")
            print("   2. Run this script again to verify")
            print("   3. Make sure soffice/libreoffice is in your PATH")

    def create_test_powerpoint(self, output_path="test_presentation.pptx"):
        """Create a simple test PowerPoint file (requires python-pptx)."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            print(f"\n📝 Creating test PowerPoint file...")

            prs = Presentation()

            # Slide 1: Title slide
            slide1 = prs.slides.add_slide(prs.slide_layouts[0])
            slide1.shapes.title.text = "Test Presentation"
            slide1.placeholders[1].text = "Created for LibreOffice testing"

            # Slide 2: Content slide with hyperlink
            slide2 = prs.slides.add_slide(prs.slide_layouts[1])
            slide2.shapes.title.text = "Slide with Content"
            content = slide2.placeholders[1]
            tf = content.text_frame
            tf.text = "This slide has various content:"
            p = tf.add_paragraph()
            p.text = "• A bullet point"
            p = tf.add_paragraph()
            p.text = "• Another point with a link"
            p = tf.add_paragraph()
            p.text = "Visit: https://example.com"

            # Save
            prs.save(output_path)
            print(f"   ✅ Created: {output_path}")
            return output_path

        except ImportError:
            print(f"\n📝 Cannot create test PowerPoint (python-pptx not installed)")
            print("   Install with: pip install python-pptx")
            return None


def main():
    """Main function."""
    parser = argparse.ArgumentParser(
        description='Check LibreOffice installation and functionality',
        formatter_class=argparse.RawDescriptionHelpFormatter
    )

    parser.add_argument(
        '--test-convert',
        help='Test conversion with a specific file',
        metavar='FILE'
    )

    parser.add_argument(
        '--create-test-ppt',
        action='store_true',
        help='Create a test PowerPoint file'
    )

    parser.add_argument(
        '--quick',
        action='store_true',
        help='Quick check only (no conversion test)'
    )

    args = parser.parse_args()

    print("🏥 LibreOffice Installation Checker")
    print("=" * 50)

    checker = LibreOfficeChecker()

    # Find LibreOffice
    checker.find_libreoffice()

    if checker.libreoffice_path:
        # Check version
        checker.check_version()

        # Test conversion (unless quick mode)
        if not args.quick:
            if args.create_test_ppt:
                test_file = checker.create_test_powerpoint()
                if test_file:
                    checker.test_conversion(test_file)
            else:
                checker.test_conversion(args.test_convert)

    # Check dependencies
    checker.check_dependencies()

    # Print summary
    checker.print_summary()

    # Return exit code
    return 0 if checker.libreoffice_path else 1


if __name__ == "__main__":
    sys.exit(main())